import streamlit as st
import nltk
from nltk.stem import WordNetLemmatizer
from nltk.corpus import stopwords
from nltk import pos_tag
import re
from gensim.models import LdaModel
from gensim import corpora
from gensim.summarization import summarize
from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import PyPDF2
from docx import Document
from pptx import Presentation
import os
from collections import Counter
import time
import random
from io import BytesIO
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from reportlab.pdfgen import canvas
import plotly.io as pio
from datetime import datetime

# ----------------------------
# Page Configuration (MUST be first Streamlit command)
# ----------------------------
st.set_page_config(
    page_title="Narrative Nexus", 
    page_icon="🤖", 
    layout="wide",
    initial_sidebar_state="collapsed"
)

# ----------------------------
# PDF REPORT GENERATION FUNCTIONS
# ----------------------------
def create_pdf_report(text, analysis_data):
    """
    Generate a comprehensive PDF report with all analysis results
    """
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, 
                           rightMargin=72, leftMargin=72,
                           topMargin=72, bottomMargin=18)
    
    # Container for PDF elements
    elements = []
    
    # Styles
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#00d4ff'),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#4ecdc4'),
        spaceAfter=12,
        spaceBefore=12,
        fontName='Helvetica-Bold'
    )
    
    subheading_style = ParagraphStyle(
        'CustomSubHeading',
        parent=styles['Heading3'],
        fontSize=12,
        textColor=colors.HexColor('#00d4ff'),
        spaceAfter=6,
        fontName='Helvetica-Bold'
    )
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['BodyText'],
        fontSize=10,
        alignment=TA_JUSTIFY,
        spaceAfter=12
    )
    
    # Title
    elements.append(Paragraph("🤖 NARRATIVE NEXUS", title_style))
    elements.append(Paragraph("AI-Powered Text Analysis Report", styles['Normal']))
    elements.append(Spacer(1, 0.3*inch))
    
    # Report metadata
    report_date = datetime.now().strftime("%B %d, %Y at %I:%M %p")
    elements.append(Paragraph(f"<b>Generated:</b> {report_date}", styles['Normal']))
    elements.append(Spacer(1, 0.2*inch))
    
    # Executive Summary
    elements.append(Paragraph("📊 EXECUTIVE SUMMARY", heading_style))
    
    summary_data = [
        ['Metric', 'Value'],
        ['Primary Domain', analysis_data.get('primary_domain', 'N/A')],
        ['Confidence Level', f"{analysis_data.get('confidence', 0):.0f}%"],
        ['Overall Sentiment', analysis_data.get('sentiment_label', 'Neutral')],
        ['Sentiment Score', f"{analysis_data.get('compound_score', 0):.3f}"],
        ['Word Count', f"{analysis_data.get('word_count', 0):,}"],
        ['Sentence Count', str(analysis_data.get('sentence_count', 0))],
        ['Vocabulary Richness', f"{analysis_data.get('vocab_richness', 0):.2f}"],
    ]
    
    summary_table = Table(summary_data, colWidths=[3*inch, 3*inch])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#00d4ff')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('FONTNAME', (0, 1), (0, -1), 'Helvetica-Bold'),
    ]))
    
    elements.append(summary_table)
    elements.append(Spacer(1, 0.3*inch))
    
    # Text Preview
    elements.append(Paragraph("📄 TEXT PREVIEW", heading_style))
    preview_text = text[:500] + ('...' if len(text) > 500 else '')
    elements.append(Paragraph(preview_text, body_style))
    elements.append(Spacer(1, 0.2*inch))
    
    # Generated Title
    elements.append(Paragraph("🏆 GENERATED TITLE", heading_style))
    elements.append(Paragraph(analysis_data.get('generated_title', 'N/A'), body_style))
    elements.append(Spacer(1, 0.2*inch))
    
    # Summary
    elements.append(Paragraph("📝 SMART SUMMARY", heading_style))
    elements.append(Paragraph(analysis_data.get('summary', 'No summary available'), body_style))
    elements.append(Spacer(1, 0.3*inch))
    
    # Domain Analysis
    elements.append(Paragraph("🌐 DOMAIN ANALYSIS", heading_style))
    
    if analysis_data.get('domain_scores'):
        domain_data = [['Domain', 'Score']]
        for domain, score in analysis_data['domain_scores'][:5]:
            domain_data.append([domain, f"{score:.2f}"])
        
        domain_table = Table(domain_data, colWidths=[4*inch, 2*inch])
        domain_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4ecdc4')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.lightgrey),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        
        elements.append(domain_table)
    
    elements.append(Spacer(1, 0.3*inch))
    
    # Sentiment Analysis Details
    elements.append(Paragraph("💭 SENTIMENT ANALYSIS", heading_style))
    
    sentiment_data = [
        ['Component', 'Score'],
        ['Positive', f"{analysis_data.get('pos_score', 0):.3f}"],
        ['Neutral', f"{analysis_data.get('neu_score', 0):.3f}"],
        ['Negative', f"{analysis_data.get('neg_score', 0):.3f}"],
        ['Compound', f"{analysis_data.get('compound_score', 0):.3f}"],
    ]
    
    sentiment_table = Table(sentiment_data, colWidths=[3*inch, 3*inch])
    sentiment_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#ff6b6b')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 11),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.lightgrey),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))
    
    elements.append(sentiment_table)
    elements.append(Spacer(1, 0.3*inch))
    
    # Page Break
    elements.append(PageBreak())
    
    # Keywords Analysis
    elements.append(Paragraph("🏷️ KEY TERMS ANALYSIS", heading_style))
    
    if analysis_data.get('top_keywords'):
        keyword_data = [['Keyword', 'Frequency']]
        for word, count in analysis_data['top_keywords'][:10]:
            keyword_data.append([word, str(count)])
        
        keyword_table = Table(keyword_data, colWidths=[4*inch, 2*inch])
        keyword_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#ffd93d')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.lightgrey),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        
        elements.append(keyword_table)
    
    elements.append(Spacer(1, 0.3*inch))
    
    # Text Characteristics
    elements.append(Paragraph("📊 TEXT CHARACTERISTICS", heading_style))
    
    char_data = [
        ['Characteristic', 'Value'],
        ['Average Words per Sentence', f"{analysis_data.get('avg_words_per_sentence', 0):.1f}"],
        ['Unique Words', f"{analysis_data.get('unique_words', 0):,}"],
        ['Content Density', f"{analysis_data.get('content_density', 0):.2f}"],
        ['Complexity Score', f"{analysis_data.get('complexity_score', 0):.1f}/100"],
    ]
    
    char_table = Table(char_data, colWidths=[3*inch, 3*inch])
    char_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4ecdc4')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 11),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.lightgrey),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))
    
    elements.append(char_table)
    elements.append(Spacer(1, 0.3*inch))
    
    # Analysis Insights
    elements.append(Paragraph("💡 KEY INSIGHTS", heading_style))
    
    insights = analysis_data.get('insights', [])
    for insight in insights:
        elements.append(Paragraph(f"• {insight}", body_style))
    
    elements.append(Spacer(1, 0.3*inch))
    
    # Optimization Suggestions
    elements.append(Paragraph("🚀 OPTIMIZATION SUGGESTIONS", heading_style))
    
    suggestions = analysis_data.get('suggestions', [])
    for suggestion in suggestions:
        elements.append(Paragraph(f"• {suggestion}", body_style))
    
    elements.append(Spacer(1, 0.3*inch))
    
    # Method Used
    elements.append(Paragraph("🔍 ANALYSIS METHOD", heading_style))
    elements.append(Paragraph(f"<b>Method Used:</b> {analysis_data.get('method_used', 'N/A')}", body_style))
    elements.append(Paragraph(f"<b>Reason:</b> {analysis_data.get('method_reason', 'N/A')}", body_style))
    
    if analysis_data.get('overlap_ratio'):
        elements.append(Paragraph(f"<b>Vocabulary Overlap:</b> {analysis_data.get('overlap_ratio', 0):.1%}", body_style))
    
    elements.append(Spacer(1, 0.5*inch))
    
    # Footer
    elements.append(Paragraph("─" * 80, styles['Normal']))
    elements.append(Paragraph("Generated by Narrative Nexus - AI-Powered Text Analysis", 
                             ParagraphStyle('Footer', parent=styles['Normal'], 
                                          fontSize=8, textColor=colors.grey, alignment=TA_CENTER)))
    
    # Build PDF
    doc.build(elements)
    
    buffer.seek(0)
    return buffer

# ----------------------------
# CUSTOM CSS AND ANIMATIONS
# ----------------------------
def load_custom_css():
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@400;700;900&family=Inter:wght@300;400;600&display=swap');
    
    /* Main background with gradient */
    .stApp {
        background: linear-gradient(135deg, #0f0f23 0%, #1a1a2e 50%, #16213e 100%);
        color: #ffffff;
    }
    
    /* Floating particles container */
    .particles {
        position: fixed;
        top: 0;
        left: 0;
        width: 100%;
        height: 100%;
        pointer-events: none;
        z-index: -1;
        overflow: hidden;
    }
    
    /* Individual floating particles */
    .particle {
        position: absolute;
        color: #00d4ff;
        font-size: 20px;
        animation: float 15s infinite linear;
        opacity: 0.6;
    }
    
    .particle:nth-child(odd) {
        color: #ff6b6b;
        animation-duration: 12s;
    }
    
    .particle:nth-child(3n) {
        color: #4ecdc4;
        animation-duration: 18s;
    }
    
    @keyframes float {
        0% {
            transform: translateY(100vh) rotate(0deg);
            opacity: 0;
        }
        10% {
            opacity: 0.6;
        }
        90% {
            opacity: 0.6;
        }
        100% {
            transform: translateY(-100px) rotate(360deg);
            opacity: 0;
        }
    }
    
    /* Main title styling */
    .main-title {
        font-family: 'Orbitron', monospace;
        font-size: 3.5rem;
        font-weight: 900;
        text-align: center;
        background: linear-gradient(45deg, #00d4ff, #ff6b6b, #4ecdc4);
        background-size: 400% 400%;
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        animation: gradient-shift 3s ease-in-out infinite;
        margin-bottom: 10px;
        text-shadow: 0 0 30px rgba(0, 212, 255, 0.5);
    }
    
    @keyframes gradient-shift {
        0%, 100% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
    }
    
    /* Subtitle styling */
    .subtitle {
        font-family: 'Inter', sans-serif;
        font-size: 1.2rem;
        text-align: center;
        color: #a0a0a0;
        margin-bottom: 30px;
        animation: pulse 2s ease-in-out infinite;
    }
    
    @keyframes pulse {
        0%, 100% { opacity: 0.8; }
        50% { opacity: 1; }
    }
    
    /* Container styling */
    .main-container {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(10px);
        border-radius: 20px;
        padding: 30px;
        margin: 20px 0;
        border: 1px solid rgba(255, 255, 255, 0.1);
        box-shadow: 0 20px 40px rgba(0, 0, 0, 0.3);
    }
    
    /* Section headers */
    .section-header {
        font-family: 'Orbitron', monospace;
        font-size: 1.8rem;
        color: #00d4ff;
        margin: 25px 0 15px 0;
        position: relative;
        display: inline-block;
    }
    
    .section-header::after {
        content: '';
        position: absolute;
        bottom: -5px;
        left: 0;
        width: 100%;
        height: 2px;
        background: linear-gradient(90deg, #00d4ff, transparent);
        animation: shimmer 2s ease-in-out infinite;
    }
    
    @keyframes shimmer {
        0% { transform: scaleX(0); }
        50% { transform: scaleX(1); }
        100% { transform: scaleX(0); }
    }
    
    /* Custom cards */
    .info-card {
        background: linear-gradient(135deg, rgba(0, 212, 255, 0.1), rgba(255, 107, 107, 0.1));
        border-radius: 15px;
        padding: 20px;
        margin: 15px 0;
        border: 1px solid rgba(0, 212, 255, 0.3);
        transition: transform 0.3s ease, box-shadow 0.3s ease;
    }
    
    .info-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 15px 30px rgba(0, 212, 255, 0.2);
    }
    
    /* Alert cards for method selection */
    .method-alert {
        background: linear-gradient(135deg, rgba(78, 205, 196, 0.15), rgba(255, 107, 107, 0.15));
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
        border-left: 4px solid #4ecdc4;
    }
    
    /* Custom buttons */
    .stButton > button {
        background: linear-gradient(45deg, #00d4ff, #4ecdc4);
        color: white;
        border: none;
        border-radius: 25px;
        padding: 12px 25px;
        font-weight: 600;
        font-family: 'Inter', sans-serif;
        transition: all 0.3s ease;
        box-shadow: 0 5px 15px rgba(0, 212, 255, 0.3);
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 10px 25px rgba(0, 212, 255, 0.5);
    }
    
    /* Download button special styling */
    .stDownloadButton > button {
        background: linear-gradient(45deg, #ff6b6b, #ff8e53);
        color: white;
        border: none;
        border-radius: 25px;
        padding: 12px 25px;
        font-weight: 600;
        font-family: 'Inter', sans-serif;
        transition: all 0.3s ease;
        box-shadow: 0 5px 15px rgba(255, 107, 107, 0.3);
    }
    
    .stDownloadButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 10px 25px rgba(255, 107, 107, 0.5);
    }
    
    /* Radio buttons */
    .stRadio > div {
        background: rgba(255, 255, 255, 0.05);
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
    }
    
    /* Text area */
    .stTextArea textarea {
        background: rgba(255, 255, 255, 0.1);
        border: 1px solid rgba(0, 212, 255, 0.3);
        border-radius: 10px;
        color: white;
    }
    
    /* File uploader */
    .stFileUploader {
        background: rgba(255, 255, 255, 0.05);
        border-radius: 10px;
        padding: 20px;
        border: 2px dashed rgba(0, 212, 255, 0.3);
        transition: border-color 0.3s ease;
    }
    
    .stFileUploader:hover {
        border-color: rgba(0, 212, 255, 0.6);
    }
    
    /* Loading spinner customization */
    .stSpinner {
        text-align: center;
        color: #00d4ff;
    }
    
    /* Hide Streamlit elements */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    /* Sidebar styling */
    .css-1d391kg {
        background: linear-gradient(180deg, #1a1a2e 0%, #16213e 100%);
    }
    
    /* Metrics styling */
    .metric-container {
        background: linear-gradient(135deg, rgba(78, 205, 196, 0.1), rgba(255, 107, 107, 0.1));
        padding: 15px;
        border-radius: 10px;
        margin: 10px 0;
        border-left: 4px solid #00d4ff;
    }
    
    /* Status indicators */
    .status-positive { color: #4ecdc4; }
    .status-negative { color: #ff6b6b; }
    .status-neutral { color: #ffd93d; }
    </style>
    """, unsafe_allow_html=True)

def create_floating_particles():
    particles_html = '<div class="particles">'
    particle_chars = ['⚡', '🔮', '✨', '💫', '🌟', '🔥', '💎', '🚀', '⭐', '💥']
    
    for i in range(20):
        char = random.choice(particle_chars)
        left = random.randint(0, 100)
        delay = random.randint(0, 15)
        particles_html += f'<div class="particle" style="left: {left}%; animation-delay: -{delay}s;">{char}</div>'
    
    particles_html += '</div>'
    st.markdown(particles_html, unsafe_allow_html=True)

def create_title_section():
    st.markdown("""
    <div class="main-container">
        <h1 class="main-title">🤖 NARRATIVE NEXUS</h1>
        <p class="subtitle">Smart Domain Detection • Adaptive Analysis • Auto-Method Selection</p>
    </div>
    """, unsafe_allow_html=True)

def create_section_header(title, icon=""):
    st.markdown(f'<div class="section-header">{icon} {title}</div>', unsafe_allow_html=True)

def create_info_card(content, card_type="default"):
    st.markdown(f'<div class="info-card">{content}</div>', unsafe_allow_html=True)

def create_method_alert(content):
    st.markdown(f'<div class="method-alert">{content}</div>', unsafe_allow_html=True)

def create_metric_card(title, value, description=""):
    st.markdown(f"""
    <div class="metric-container">
        <h3 style="margin: 0; color: #00d4ff;">{title}</h3>
        <h2 style="margin: 10px 0; color: white;">{value}</h2>
        <p style="margin: 0; color: #a0a0a0; font-size: 0.9rem;">{description}</p>
    </div>
    """, unsafe_allow_html=True)

# ----------------------------
# NLTK setup
# ----------------------------
try:
    @st.cache_data
    def setup_nltk():
        nltk.download('punkt', quiet=True)
        nltk.download('wordnet', quiet=True)
        nltk.download('omw-1.4', quiet=True)
        nltk.download('stopwords', quiet=True)
        nltk.download('averaged_perceptron_tagger', quiet=True)
        return True
except AttributeError:
    @st.cache(allow_output_mutation=True)
    def setup_nltk():
        nltk.download('punkt', quiet=True)
        nltk.download('wordnet', quiet=True)
        nltk.download('omw-1.4', quiet=True)
        nltk.download('stopwords', quiet=True)
        nltk.download('averaged_perceptron_tagger', quiet=True)
        return True

setup_nltk()

lemmatizer = WordNetLemmatizer()
stop_words = set(stopwords.words('english'))
analyzer = SentimentIntensityAnalyzer()

# ----------------------------
# Load trained LDA + dictionary
# ----------------------------
MODEL_DIR = "models"
DICT_PATH = os.path.join(MODEL_DIR, "dictionary.dict")
LDA_PATH = os.path.join(MODEL_DIR, "lda.model")

lda_ready = os.path.exists(DICT_PATH) and os.path.exists(LDA_PATH)
if lda_ready:
    dictionary = corpora.Dictionary.load(DICT_PATH)
    lda = LdaModel.load(LDA_PATH)

# ----------------------------
# Domain Detection Functions
# ----------------------------
def detect_text_domain(text, noun_tokens):
    """
    Analyze text to determine its domain/category based on vocabulary and patterns.
    """
    text_lower = text.lower()
    words_set = set(nltk.word_tokenize(text_lower))
    
    # Domain indicator keywords with weights
    domain_patterns = {
        'News & Current Affairs': {
            'keywords': ['breaking', 'reported', 'according', 'sources', 'reuters', 'ap', 'cnn', 'bbc', 
                        'journalist', 'correspondent', 'statement', 'official', 'government', 'president',
                        'minister', 'election', 'political', 'parliament', 'congress', 'senate'],
            'weight': 1.0
        },
        'Product Reviews': {
            'keywords': ['review', 'rating', 'stars', 'recommend', 'purchase', 'bought', 'quality',
                        'price', 'value', 'money', 'worth', 'excellent', 'terrible', 'amazing',
                        'disappointed', 'satisfied', 'customer', 'product', 'item', 'delivery'],
            'weight': 1.2
        },
        'Entertainment & Media': {
            'keywords': ['movie', 'film', 'actor', 'actress', 'director', 'cinema', 'theater',
                        'music', 'song', 'album', 'artist', 'band', 'concert', 'performance',
                        'show', 'episode', 'season', 'character', 'plot', 'story', 'entertainment'],
            'weight': 1.1
        },
        'Technology & Science': {
            'keywords': ['technology', 'tech', 'software', 'hardware', 'algorithm', 'data',
                        'research', 'study', 'experiment', 'science', 'scientific', 'innovation',
                        'development', 'programming', 'computer', 'digital', 'artificial', 'intelligence'],
            'weight': 1.0
        },
        'Business & Finance': {
            'keywords': ['business', 'company', 'corporation', 'market', 'stock', 'investment',
                        'profit', 'revenue', 'sales', 'financial', 'economy', 'economic',
                        'bank', 'banking', 'finance', 'money', 'dollar', 'quarter', 'earnings'],
            'weight': 1.0
        },
        'Health & Medical': {
            'keywords': ['health', 'medical', 'doctor', 'patient', 'treatment', 'disease',
                        'medicine', 'therapy', 'hospital', 'clinic', 'symptoms', 'diagnosis',
                        'pharmaceutical', 'drug', 'vaccine', 'healthcare', 'wellness'],
            'weight': 1.1
        },
        'Sports': {
            'keywords': ['game', 'team', 'player', 'match', 'season', 'championship', 'tournament',
                        'score', 'goal', 'win', 'lost', 'victory', 'defeat', 'coach', 'league',
                        'football', 'basketball', 'baseball', 'soccer', 'sports'],
            'weight': 1.2
        },
        'Education & Academic': {
            'keywords': ['university', 'college', 'school', 'student', 'teacher', 'professor',
                        'education', 'academic', 'learning', 'course', 'curriculum', 'degree',
                        'research', 'study', 'knowledge', 'scholarly', 'publication'],
            'weight': 1.0
        },
        'Travel & Tourism': {
            'keywords': ['travel', 'trip', 'vacation', 'hotel', 'restaurant', 'tourist', 'tourism',
                        'destination', 'flight', 'airport', 'booking', 'reservation', 'sightseeing',
                        'culture', 'local', 'experience', 'adventure'],
            'weight': 1.1
        },
        'Lifestyle & Personal': {
            'keywords': ['life', 'personal', 'experience', 'feeling', 'think', 'believe',
                        'opinion', 'lifestyle', 'daily', 'routine', 'family', 'friends',
                        'relationship', 'love', 'happiness', 'advice', 'tips'],
            'weight': 0.9
        }
    }
    
    domain_scores = {}
    
    for domain, config in domain_patterns.items():
        score = 0
        keywords = config['keywords']
        weight = config['weight']
        
        # Count keyword matches
        for keyword in keywords:
            if keyword in text_lower:
                score += text_lower.count(keyword) * weight
        
        # Bonus for exact phrase matches
        phrases = {
            'News & Current Affairs': ['breaking news', 'press release', 'according to sources'],
            'Product Reviews': ['five stars', 'would recommend', 'great value'],
            'Entertainment & Media': ['box office', 'red carpet', 'behind the scenes'],
            'Technology & Science': ['machine learning', 'artificial intelligence', 'peer review'],
            'Business & Finance': ['quarterly earnings', 'market cap', 'stock price'],
            'Sports': ['home team', 'final score', 'world championship']
        }
        
        if domain in phrases:
            for phrase in phrases[domain]:
                if phrase in text_lower:
                    score += 3 * weight
        
        domain_scores[domain] = score
    
    # Normalize scores by text length
    text_length = len(nltk.word_tokenize(text))
    if text_length > 0:
        for domain in domain_scores:
            domain_scores[domain] = domain_scores[domain] / text_length * 100
    
    # Sort domains by score
    sorted_domains = sorted(domain_scores.items(), key=lambda x: x[1], reverse=True)
    
    # Calculate confidence based on score difference
    if len(sorted_domains) > 1:
        top_score = sorted_domains[0][1]
        second_score = sorted_domains[1][1]
        confidence = min(95, max(60, (top_score - second_score) * 10 + 60))
    else:
        confidence = 60
    
    # Return top 5 domains with scores
    return sorted_domains[:5], confidence

def calculate_vocabulary_overlap(noun_tokens, dictionary, min_overlap_ratio=0.3, min_words_in_vocab=10):
    """
    Calculate vocabulary overlap between current text and LDA model's dictionary.
    """
    if not dictionary or not noun_tokens:
        return 0.0, "top_nouns", "No vocabulary or dictionary available"
    
    bow = dictionary.doc2bow(noun_tokens)
    words_in_vocab = len(bow)
    total_unique_words = len(set(noun_tokens))
    
    if total_unique_words == 0:
        return 0.0, "top_nouns", "No words found"
    
    overlap_ratio = words_in_vocab / total_unique_words
    
    if words_in_vocab < min_words_in_vocab:
        return overlap_ratio, "top_nouns", f"Too few vocabulary matches ({words_in_vocab} < {min_words_in_vocab})"
    elif overlap_ratio < min_overlap_ratio:
        return overlap_ratio, "top_nouns", f"Low vocabulary overlap ({overlap_ratio:.1%} < {min_overlap_ratio:.1%})"
    else:
        return overlap_ratio, "lda", f"Good vocabulary overlap ({overlap_ratio:.1%})"

def get_lda_domain_indicators(dictionary):
    """
    Analyze the LDA model's vocabulary to determine what domain it was trained on.
    """
    if not dictionary:
        return "Unknown domain"
    
    domain_indicators = {
        'news_business': ['reuters', 'company', 'inc', 'corp', 'said', 'wednesday', 'thursday', 'tuesday', 'monday', 'friday'],
        'academic': ['research', 'study', 'analysis', 'theory', 'method', 'paper', 'journal', 'university'],
        'social_media': ['tweet', 'like', 'share', 'follow', 'hashtag', 'post', 'comment'],
        'literature': ['character', 'story', 'novel', 'author', 'book', 'chapter', 'plot'],
        'technical': ['system', 'data', 'algorithm', 'process', 'method', 'implementation', 'framework']
    }
    
    vocab_words = set(dictionary.token2id.keys())
    domain_scores = {}
    
    for domain, indicators in domain_indicators.items():
        score = sum(1 for word in indicators if word in vocab_words)
        domain_scores[domain] = score
    
    if max(domain_scores.values()) > 0:
        likely_domain = max(domain_scores, key=domain_scores.get)
        return likely_domain.replace('_', '/').title()
    else:
        return "General/Mixed domain"

# ----------------------------
# Preprocessing functions
# ----------------------------
def clean_text(text):
    text = text.lower()
    text = re.sub(r"http\S+|www\S+|https\S+", "", text)
    text = re.sub(r"<.*?>", "", text)
    text = re.sub(r"[^a-z\s]", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text

def remove_stopwords(text):
    toks = nltk.word_tokenize(text)
    toks = [w for w in toks if w not in stop_words]
    return " ".join(toks)

def lemmatize_text(text):
    toks = nltk.word_tokenize(text)
    toks = [lemmatizer.lemmatize(w) for w in toks]
    return " ".join(toks)

def nouns_only_tokens(text):
    toks = nltk.word_tokenize(text)
    return [w for w, p in pos_tag(toks) if p.startswith("NN") and len(w) > 1]

def top_nouns(tokens, k=10):
    c = Counter(tokens)
    return c.most_common(k), sum(c.values())

# ----------------------------
# Enhanced Plotting Functions
# ----------------------------
def create_enhanced_bar_chart(x, y, title, x_label, y_label, color_scheme='viridis'):
    fig = go.Figure(data=[
        go.Bar(
            x=x, y=y,
            marker=dict(
                color=y,
                colorscale=color_scheme,
                showscale=True,
                colorbar=dict(thickness=15, len=0.7)
            ),
            text=[f'{val:.2f}' if isinstance(val, float) else str(val) for val in y],
            textposition='outside',
            hovertemplate='<b>%{x}</b><br>Value: %{y}<extra></extra>'
        )
    ])
    
    fig.update_layout(
        title=dict(
            text=title,
            font=dict(size=20, color='white', family='Orbitron'),
            x=0.5
        ),
        xaxis_title=x_label,
        yaxis_title=y_label,
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        font=dict(color='white'),
        xaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
        yaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
        margin=dict(t=60, b=60, l=60, r=60)
    )
    
    return fig

def create_sentiment_gauge(compound_score):
    fig = go.Figure(go.Indicator(
        mode = "gauge+number+delta",
        value = compound_score,
        domain = {'x': [0, 1], 'y': [0, 1]},
        title = {'text': "Sentiment Score", 'font': {'color': 'white', 'size': 20}},
        delta = {'reference': 0, 'increasing': {'color': "RebeccaPurple"}},
        gauge = {
            'axis': {'range': [-1, 1], 'tickcolor': "white"},
            'bar': {'color': "darkblue"},
            'steps': [
                {'range': [-1, -0.05], 'color': "#ff6b6b"},
                {'range': [-0.05, 0.05], 'color': "#ffd93d"},
                {'range': [0.05, 1], 'color': "#4ecdc4"}
            ],
            'threshold': {
                'line': {'color': "red", 'width': 4},
                'thickness': 0.75,
                'value': 0.9
            }
        }
    ))
    
    fig.update_layout(
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        font={'color': "white", 'size': 12},
        height=400
    )
    
    return fig

def create_domain_radar_chart(domain_scores, confidence):
    """Create a radar chart showing domain probabilities"""
    domains = [domain for domain, _ in domain_scores]
    scores = [score for _, score in domain_scores]
    
    fig = go.Figure()
    
    fig.add_trace(go.Scatterpolar(
        r=scores,
        theta=domains,
        fill='toself',
        name='Domain Probability',
        line=dict(color='#00d4ff', width=2),
        fillcolor='rgba(0, 212, 255, 0.3)'
    ))
    
    fig.update_layout(
        polar=dict(
            radialaxis=dict(
                visible=True,
                range=[0, max(scores) * 1.1] if scores else [0, 1],
                gridcolor='rgba(255,255,255,0.2)',
                tickcolor='white'
            ),
            angularaxis=dict(
                gridcolor='rgba(255,255,255,0.2)',
                tickcolor='white'
            )
        ),
        showlegend=True,
        title=dict(
            text=f"🎯 Domain Analysis - {confidence:.0f}% Confidence",
            font=dict(size=20, color='white', family='Orbitron'),
            x=0.5
        ),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        font=dict(color='white'),
        height=500
    )
    
    return fig

# ----------------------------
# Text Analysis Functions
# ----------------------------
def analyze_text_characteristics(text, noun_tokens):
    """Analyze various characteristics of the text"""
    words = nltk.word_tokenize(text)
    sentences = nltk.sent_tokenize(text)
    
    word_count = len(words)
    sentence_count = len(sentences)
    avg_words_per_sentence = word_count / sentence_count if sentence_count > 0 else 0
    
    unique_words = len(set(word.lower() for word in words if word.isalpha()))
    vocab_richness = unique_words / word_count if word_count > 0 else 0
    
    if sentence_count > 0 and word_count > 0:
        avg_sentence_length = word_count / sentence_count
        complexity_score = min(100, max(0, (avg_sentence_length - 10) * 5 + 50))
    else:
        complexity_score = 50
    
    content_density = len(noun_tokens) / word_count if word_count > 0 else 0
    
    return {
        'word_count': word_count,
        'sentence_count': sentence_count,
        'avg_words_per_sentence': avg_words_per_sentence,
        'vocab_richness': vocab_richness,
        'complexity_score': complexity_score,
        'content_density': content_density,
        'unique_words': unique_words
    }

# ----------------------------
# Main Application Logic
# ----------------------------

# Load custom CSS and create floating particles
load_custom_css()
create_floating_particles()

# Title section
create_title_section()

# Main interface
col1, col2 = st.columns([2, 1])

with col1:
    create_section_header("🔍 INPUT CONFIGURATION", "⚙️")
    
    input_type = st.radio(
        'Select Input Method:', 
        ('✏️ Enter Text Manually', '📂 Upload Document'), 
        horizontal=True
    )
    
    # Analysis method selection
    create_section_header("🧠 ANALYSIS METHOD", "🔬")
    topic_method = st.radio(
        'Choose Analysis Method:',
        ('🤖 Smart Auto-Select (Recommended)', '🏷️ Top Keywords Only', '🧠 LDA Only'),
        horizontal=False,
        help="Smart Auto-Select analyzes vocabulary overlap to choose the best method automatically"
    )

with col2:
    create_section_header("📊 ANALYSIS INFO", "📈")
    
    if lda_ready:
        domain = get_lda_domain_indicators(dictionary)
        vocab_size = len(dictionary.token2id)
        st.markdown(f"""
        <div class="info-card">
            <h4 style="color: #4ecdc4; margin: 0;">✅ LDA Model Ready</h4>
            <p style="color: #a0a0a0; margin: 5px 0;"><strong>Domain:</strong> {domain}</p>
            <p style="color: #a0a0a0; margin: 5px 0;"><strong>Vocabulary:</strong> {vocab_size:,} words</p>
            <p style="color: #a0a0a0; margin: 5px 0;"><strong>Topics:</strong> {lda.num_topics}</p>
        </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown("""
        <div class="info-card">
            <h4 style="color: #ff6b6b; margin: 0;">❌ LDA Model Not Found</h4>
            <p style="color: #a0a0a0; margin: 5px 0;">Will use Domain Detection analysis</p>
        </div>
        """, unsafe_allow_html=True)

# Text input section
text = ''
create_section_header("📄 TEXT INPUT", "📝")

if input_type == '✏️ Enter Text Manually':
    text = st.text_area(
        '✨ Enter your text here for AI analysis:', 
        height=200,
        placeholder="Paste your text here... The AI will automatically detect the domain and select the best analysis method!"
    )
else:
    uploaded_file = st.file_uploader(
        '🚀 Upload your document:', 
        type=['pdf', 'txt', 'docx', 'pptx'],
        help="Supported formats: PDF, TXT, DOCX, PPTX"
    )
    
    if uploaded_file is not None:
        with st.spinner('🔄 Processing your document...'):
            file_type = uploaded_file.type
            try:
                if file_type == 'text/plain':
                    text = uploaded_file.read().decode('utf-8', errors='ignore')
                elif file_type == 'application/pdf':
                    pdf_reader = PyPDF2.PdfReader(uploaded_file)
                    text = "".join([page.extract_text() or "" for page in pdf_reader.pages])
                elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    doc = Document(uploaded_file)
                    text = "\n".join([para.text for para in doc.paragraphs])
                elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                    ppt = Presentation(uploaded_file)
                    text = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
                else:
                    st.error('❌ Unsupported file type.')
            except Exception as e:
                st.error(f'❌ Error reading file: {e}')

# Analysis section
if text:
    # Text preview
    create_section_header("👀 TEXT PREVIEW", "📖")
    preview_text = text[:500] + ('...' if len(text) > 500 else '')
    create_info_card(f"<pre style='white-space: pre-wrap; font-size: 0.9rem;'>{preview_text}</pre>")
    
    # Processing with spinner
    with st.spinner('🧠 AI is analyzing your text... Please wait'):
        time.sleep(1)
        cleaned = clean_text(text)
        no_stop = remove_stopwords(cleaned)
        lemmad = lemmatize_text(no_stop)
        noun_tokens = nouns_only_tokens(lemmad)
    
    # Smart method selection logic
    selected_method = "top_nouns"
    overlap_ratio = 0.0
    selection_reason = ""
    
    if topic_method == '🤖 Smart Auto-Select (Recommended)' and lda_ready:
        overlap_ratio, selected_method, selection_reason = calculate_vocabulary_overlap(noun_tokens, dictionary)
        
        create_section_header("🎯 SMART METHOD SELECTION", "🤖")
        if selected_method == "lda":
            create_method_alert(f"""
                <strong>🧠 Using LDA Analysis</strong><br>
                <strong>Vocabulary Overlap:</strong> {overlap_ratio:.1%}<br>
                <strong>Reason:</strong> {selection_reason}<br>
                <strong>✅ Good match with training domain</strong>
            """)
        else:
            create_method_alert(f"""
                <strong>🏷️ Using Domain Detection</strong><br>
                <strong>Vocabulary Overlap:</strong> {overlap_ratio:.1%}<br>
                <strong>Reason:</strong> {selection_reason}<br>
                <strong>⚠️ Text domain differs from LDA training data</strong>
            """)
    elif topic_method == '🏷️ Top Keywords Only':
        selected_method = "top_nouns"
        create_method_alert("🏷️ <strong>Using Domain Detection</strong> (User Selected)")
    elif topic_method == '🧠 LDA Only':
        if lda_ready:
            selected_method = "lda"
            create_method_alert("🧠 <strong>Using LDA Analysis</strong> (User Selected)")
        else:
            selected_method = "top_nouns"
            create_method_alert("🏷️ <strong>Fallback to Domain Detection</strong> (LDA model not available)")
    
    # Domain Analysis
    domain_scores, confidence = detect_text_domain(text, noun_tokens)
    text_characteristics = analyze_text_characteristics(text, noun_tokens)
    
    # Results in columns
    col1, col2 = st.columns(2)
    
    # Initialize variables for PDF report
    generated_title = ""
    summary_text = ""
    
    with col1:
        # Title Generation
        create_section_header("🏆 GENERATED TITLE", "📰")
        if noun_tokens:
            top1, _ = top_nouns(noun_tokens, k=1)
            title_word = top1[0][0].title()
            generated_title = f"📚 {title_word}"
            create_metric_card("Generated Title", generated_title, "Based on most frequent key noun")
        else:
            fallback = " ".join(nltk.word_tokenize(lemmad)[:10]).capitalize() + "..."
            generated_title = f"📚 {fallback}"
            create_metric_card("Generated Title", generated_title, "Fallback title generated")
        
        # Summary
        create_section_header("📝 SMART SUMMARY", "✨")
        try:
            with st.spinner('✨ Generating intelligent summary...'):
                summary = None
                for ratio in [0.3, 0.2, 0.1]:
                    try:
                        summary = summarize(text, ratio=ratio)
                        if summary and len(summary.strip()) > 50:
                            break
                    except:
                        continue
                
                if not summary:
                    sentences = nltk.sent_tokenize(text)
                    if len(sentences) > 2:
                        summary = '. '.join(sentences[:2]) + '.'
                    else:
                        summary = text[:300] + ('...' if len(text) > 300 else '')
                
                summary_text = summary
                create_info_card(f"<div style='font-size: 1.1rem; line-height: 1.6;'>{summary}</div>")
        except Exception as e:
            sentences = nltk.sent_tokenize(text)
            if len(sentences) > 1:
                summary = '. '.join(sentences[:2]) + '.'
            else:
                summary = text[:300] + ('...' if len(text) > 300 else '')
            summary_text = summary
            create_info_card(f"<div style='font-size: 1.1rem; line-height: 1.6;'>{summary}</div>")
    
    with col2:
        # Sentiment Analysis
        create_section_header("💭 SENTIMENT ANALYSIS", "🎭")
        sentiment = analyzer.polarity_scores(text)
        compound = sentiment['compound']
        
        if compound >= 0.05:
            mood = 'Positive 😊'
            mood_color = 'status-positive'
        elif compound <= -0.05:
            mood = 'Negative 😞'
            mood_color = 'status-negative'
        else:
            mood = 'Neutral 😐'
            mood_color = 'status-neutral'
        
        create_metric_card(
            "Overall Sentiment", 
            f"<span class='{mood_color}'>{mood}</span>",
            f"Confidence Score: {abs(compound):.2f}"
        )
        
        # Sentiment Gauge
        fig_gauge = create_sentiment_gauge(compound)
        st.plotly_chart(fig_gauge, use_container_width=True)
    
    # Domain Analysis (Full Width)
    create_section_header("🌐 DOMAIN DETECTION", "📊")
    
    if domain_scores:
        primary_domain = domain_scores[0][0]
        primary_score = domain_scores[0][1]
        
        create_metric_card(
            "Detected Domain", 
            f"🎯 {primary_domain}",
            f"Confidence: {confidence:.0f}% | Score: {primary_score:.2f}"
        )
        
        # Domain radar chart
        if len(domain_scores) > 1:
            fig_radar = create_domain_radar_chart(domain_scores, confidence)
            st.plotly_chart(fig_radar, use_container_width=True)
        
        # Top domains bar chart
        domains = [domain for domain, _ in domain_scores]
        scores = [score for _, score in domain_scores]
        
        fig_domains = create_enhanced_bar_chart(
            domains, scores,
            '🎯 Domain Probability Scores',
            'Domains', 'Probability Score',
            'plasma'
        )
        st.plotly_chart(fig_domains, use_container_width=True)
    else:
        create_info_card("⚠️ Unable to detect domain from the provided text.")
    
    # Keyword Analysis
    top_keywords = []
    if selected_method == "top_nouns":
        create_section_header("🏷️ KEY TERMS ANALYSIS", "🔑")
        if noun_tokens:
            top_items, total = top_nouns(noun_tokens, k=10)
            top_keywords = top_items
            words = [w for w, c in top_items]
            counts = [c for w, c in top_items]
            
            create_metric_card(
                "Most Frequent Term", 
                f"🏅 {words[0]}", 
                f"Appears {counts[0]} times ({counts[0]/total:.1%} of all key terms)"
            )
            
            fig_noun = create_enhanced_bar_chart(
                words, counts,
                '🏷️ Top Keywords Analysis',
                'Keywords', 'Frequency',
                'viridis'
            )
            st.plotly_chart(fig_noun, use_container_width=True)
        else:
            create_info_card("⚠️ No significant terms found after preprocessing.")
    
    else:  # LDA Analysis
        if noun_tokens:
            top_items, total = top_nouns(noun_tokens, k=10)
            top_keywords = top_items
    
    # Text Characteristics Analysis
    create_section_header("📊 TEXT CHARACTERISTICS", "📈")
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        create_metric_card("Word Count", f"{text_characteristics['word_count']:,}", "Total words in text")
    with col2:
        create_metric_card("Sentences", f"{text_characteristics['sentence_count']}", "Number of sentences")
    with col3:
        create_metric_card("Vocabulary Richness", f"{text_characteristics['vocab_richness']:.2f}", "Unique words / Total words")
    with col4:
        create_metric_card("Content Density", f"{text_characteristics['content_density']:.2f}", "Key terms / Total words")
    
    # Additional Sentiment Details
    create_section_header("📈 DETAILED SENTIMENT BREAKDOWN", "📉")
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        create_metric_card("Positive", f"{sentiment['pos']:.2f}", "Joy, happiness, optimism")
    with col2:
        create_metric_card("Neutral", f"{sentiment['neu']:.2f}", "Objective, factual tone")
    with col3:
        create_metric_card("Negative", f"{sentiment['neg']:.2f}", "Sadness, anger, criticism")
    with col4:
        create_metric_card("Compound", f"{sentiment['compound']:.2f}", "Overall sentiment score")
    
    # Sentiment breakdown chart
    scores_list = [sentiment['pos'], sentiment['neu'], sentiment['neg']]
    labels = ['Positive 😊', 'Neutral 😐', 'Negative 😞']
    colors_list = ['#4ecdc4', '#ffd93d', '#ff6b6b']
    
    fig_sentiment = go.Figure(data=[
        go.Bar(x=labels, y=scores_list, marker_color=colors_list,
              text=[f'{score:.2f}' for score in scores_list],
              textposition='outside')
    ])
    
    fig_sentiment.update_layout(
        title=dict(text='🎭 Sentiment Component Analysis', 
                  font=dict(size=20, color='white', family='Orbitron'), x=0.5),
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        font=dict(color='white'),
        xaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
        yaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
        showlegend=False
    )
    
    st.plotly_chart(fig_sentiment, use_container_width=True)
    
    # Generate insights and suggestions for PDF
    insights = []
    suggestions = []
    
    # Domain-based insights
    if domain_scores:
        primary_domain = domain_scores[0][0]
        if confidence > 80:
            insights.append(f"Strong {primary_domain.lower()} content detected")
        elif confidence > 60:
            insights.append(f"Likely {primary_domain.lower()} content")
        else:
            insights.append("Mixed domain content detected")
    
    # Sentiment insights
    if abs(compound) > 0.3:
        insights.append(f"{'Strong positive' if compound > 0 else 'Strong negative'} emotional tone detected")
    else:
        insights.append("Balanced, neutral tone")
    
    # Text quality insights
    if text_characteristics['vocab_richness'] > 0.6:
        insights.append("Rich vocabulary usage")
    elif text_characteristics['vocab_richness'] < 0.4:
        insights.append("Consider diversifying vocabulary")
    
    # Content density insights
    if text_characteristics['content_density'] > 0.3:
        insights.append("High information density")
    else:
        insights.append("Conversational style content")
    
    # Generate suggestions
    if text_characteristics['avg_words_per_sentence'] > 25:
        suggestions.append("Consider shorter sentences for better readability")
    elif text_characteristics['avg_words_per_sentence'] < 10:
        suggestions.append("Consider combining short sentences for better flow")
    else:
        suggestions.append("Good sentence length balance")
    
    if domain_scores:
        domain_suggestions = {
            'News & Current Affairs': 'Focus on timeliness and credibility',
            'Product Reviews': 'Emphasize pros/cons and value proposition',
            'Entertainment & Media': 'Highlight emotional engagement',
            'Technology & Science': 'Ensure technical accuracy',
            'Business & Finance': 'Include relevant metrics and data',
            'Health & Medical': 'Prioritize accuracy and safety',
        }