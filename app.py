import streamlit as st
import nltk
from nltk.stem import WordNetLemmatizer
from nltk.corpus import stopwords
from nltk import pos_tag
import re
from gensim.models import LdaModel
from gensim import corpora
from gensim.summarization import summarize
from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import PyPDF2
from docx import Document
from pptx import Presentation
import os
from collections import Counter
import time
import random
from io import BytesIO
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet


# ----------------------------
# Page Configuration (MUST be first Streamlit command)
# ----------------------------
st.set_page_config(
    page_title="Narrative Nexus", 
    page_icon="🤖", 
    layout="wide",
    initial_sidebar_state="collapsed"
)

# ----------------------------
# CUSTOM CSS AND ANIMATIONS
# ----------------------------
def load_custom_css():
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@400;700;900&family=Inter:wght@300;400;600&display=swap');
    
    /* Main background with gradient */
    .stApp {
        background: linear-gradient(135deg, #0f0f23 0%, #1a1a2e 50%, #16213e 100%);
        color: #ffffff;
    }
    
    /* Floating particles container */
    .particles {
        position: fixed;
        top: 0;
        left: 0;
        width: 100%;
        height: 100%;
        pointer-events: none;
        z-index: -1;
        overflow: hidden;
    }
    
    /* Individual floating particles */
    .particle {
        position: absolute;
        color: #00d4ff;
        font-size: 20px;
        animation: float 15s infinite linear;
        opacity: 0.6;
    }
    
    .particle:nth-child(odd) {
        color: #ff6b6b;
        animation-duration: 12s;
    }
    
    .particle:nth-child(3n) {
        color: #4ecdc4;
        animation-duration: 18s;
    }
    
    @keyframes float {
        0% {
            transform: translateY(100vh) rotate(0deg);
            opacity: 0;
        }
        10% {
            opacity: 0.6;
        }
        90% {
            opacity: 0.6;
        }
        100% {
            transform: translateY(-100px) rotate(360deg);
            opacity: 0;
        }
    }
    
    /* Main title styling */
    .main-title {
        font-family: 'Orbitron', monospace;
        font-size: 3.5rem;
        font-weight: 900;
        text-align: center;
        background: linear-gradient(45deg, #00d4ff, #ff6b6b, #4ecdc4);
        background-size: 400% 400%;
        -webkit-background-clip: text;
        
        animation: gradient-shift 3s ease-in-out infinite;
        margin-bottom: 10px;
        text-shadow: 0 0 30px rgba(0, 212, 255, 0.5);
    }
    
    @keyframes gradient-shift {
        0%, 100% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
    }
    
    /* Subtitle styling */
    .subtitle {
        font-family: 'Inter', sans-serif;
        font-size: 1.2rem;
        text-align: center;
        color: #a0a0a0;
        margin-bottom: 30px;
        animation: pulse 2s ease-in-out infinite;
    }
    
    @keyframes pulse {
        0%, 100% { opacity: 0.8; }
        50% { opacity: 1; }
    }
    
    /* Container styling */
    .main-container {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(10px);
        border-radius: 20px;
        padding: 30px;
        margin: 20px 0;
        border: 1px solid rgba(255, 255, 255, 0.1);
        box-shadow: 0 20px 40px rgba(0, 0, 0, 0.3);
    }
    
    /* Section headers */
    .section-header {
        font-family: 'Orbitron', monospace;
        font-size: 1.8rem;
        color: #00d4ff;
        margin: 25px 0 15px 0;
        position: relative;
        display: inline-block;
    }
    
    .section-header::after {
        content: '';
        position: absolute;
        bottom: -5px;
        left: 0;
        width: 100%;
        height: 2px;
        background: linear-gradient(90deg, #00d4ff, transparent);
        animation: shimmer 2s ease-in-out infinite;
    }
    
    @keyframes shimmer {
        0% { transform: scaleX(0); }
        50% { transform: scaleX(1); }
        100% { transform: scaleX(0); }
    }
    
    /* Custom cards */
    .info-card {
        background: linear-gradient(135deg, rgba(0, 212, 255, 0.1), rgba(255, 107, 107, 0.1));
        border-radius: 15px;
        padding: 20px;
        margin: 15px 0;
        border: 1px solid rgba(0, 212, 255, 0.3);
        transition: transform 0.3s ease, box-shadow 0.3s ease;
    }
    
    .info-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 15px 30px rgba(0, 212, 255, 0.2);
    }
    
    /* Alert cards for method selection */
    .method-alert {
        background: linear-gradient(135deg, rgba(78, 205, 196, 0.15), rgba(255, 107, 107, 0.15));
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
        border-left: 4px solid #4ecdc4;
    }
    
    /* Custom buttons */
    .stButton > button {
        background: linear-gradient(45deg, #00d4ff, #4ecdc4);
        color: white;
        border: none;
        border-radius: 25px;
        padding: 12px 25px;
        font-weight: 600;
        font-family: 'Inter', sans-serif;
        transition: all 0.3s ease;
        box-shadow: 0 5px 15px rgba(0, 212, 255, 0.3);
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 10px 25px rgba(0, 212, 255, 0.5);
    }
    
    /* Radio buttons */
    .stRadio > div {
        background: rgba(255, 255, 255, 0.05);
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
    }
    
    /* Text area */
    .stTextArea textarea {
        background: rgba(255, 255, 255, 0.1);
        border: 1px solid rgba(0, 212, 255, 0.3);
        border-radius: 10px;
        color: white;
    }
    
    /* File uploader */
    .stFileUploader {
        background: rgba(255, 255, 255, 0.05);
        border-radius: 10px;
        padding: 20px;
        border: 2px dashed rgba(0, 212, 255, 0.3);
        transition: border-color 0.3s ease;
    }
    
    .stFileUploader:hover {
        border-color: rgba(0, 212, 255, 0.6);
    }
    
    /* Loading spinner customization */
    .stSpinner {
        text-align: center;
        color: #00d4ff;
    }
    
    /* Hide Streamlit elements */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    /* Sidebar styling */
    .css-1d391kg {
        background: linear-gradient(180deg, #1a1a2e 0%, #16213e 100%);
    }
    
    /* Metrics styling */
    .metric-container {
        background: linear-gradient(135deg, rgba(78, 205, 196, 0.1), rgba(255, 107, 107, 0.1));
        padding: 15px;
        border-radius: 10px;
        margin: 10px 0;
        border-left: 4px solid #00d4ff;
    }
    
    /* Status indicators */
    .status-positive { color: #4ecdc4; }
    .status-negative { color: #ff6b6b; }
    .status-neutral { color: #ffd93d; }
    </style>
    """, unsafe_allow_html=True)

def create_floating_particles():
    particles_html = '<div class="particles">'
    particle_chars = ['⚡', '🔮', '✨', '💫', '🌟', '🔥', '💎', '🚀', '⭐', '💥']
    
    for i in range(20):
        char = random.choice(particle_chars)
        left = random.randint(0, 100)
        delay = random.randint(0, 15)
        particles_html += f'<div class="particle" style="left: {left}%; animation-delay: -{delay}s;">{char}</div>'
    
    particles_html += '</div>'
    st.markdown(particles_html, unsafe_allow_html=True)

def create_title_section():
    st.markdown("""
    <div class="main-container">
        <h1 class="main-title">🤖 NARRATIVE NEXUS</h1>
        <p class="subtitle">Smart Domain Detection • Adaptive Analysis • Auto-Method Selection</p>
    </div>
    """, unsafe_allow_html=True)

def create_section_header(title, icon=""):
    st.markdown(f'<div class="section-header">{icon} {title}</div>', unsafe_allow_html=True)

def create_info_card(content, card_type="default"):
    st.markdown(f'<div class="info-card">{content}</div>', unsafe_allow_html=True)

def create_method_alert(content):
    st.markdown(f'<div class="method-alert">{content}</div>', unsafe_allow_html=True)

def create_metric_card(title, value, description=""):
    st.markdown(f"""
    <div class="metric-container">
        <h3 style="margin: 0; color: #00d4ff;">{title}</h3>
        <h2 style="margin: 10px 0; color: white;">{value}</h2>
        <p style="margin: 0; color: #a0a0a0; font-size: 0.9rem;">{description}</p>
    </div>
    """, unsafe_allow_html=True)

# ----------------------------
# NLTK setup
# ----------------------------
try:
    # Try the new caching method first
    @st.cache_data
    def setup_nltk():
        nltk.download('punkt', quiet=True)
        nltk.download('wordnet', quiet=True)
        nltk.download('omw-1.4', quiet=True)
        nltk.download('stopwords', quiet=True)
        nltk.download('averaged_perceptron_tagger', quiet=True)
        return True
except AttributeError:
    # Fallback to old caching method
    @st.cache(allow_output_mutation=True)
    def setup_nltk():
        nltk.download('punkt', quiet=True)
        nltk.download('wordnet', quiet=True)
        nltk.download('omw-1.4', quiet=True)
        nltk.download('stopwords', quiet=True)
        nltk.download('averaged_perceptron_tagger', quiet=True)
        return True

setup_nltk()

lemmatizer = WordNetLemmatizer()
stop_words = set(stopwords.words('english'))
analyzer = SentimentIntensityAnalyzer()

# ----------------------------
# Load trained LDA + dictionary
# ----------------------------
MODEL_DIR = "models"
DICT_PATH = os.path.join(MODEL_DIR, "dictionary.dict")
LDA_PATH = os.path.join(MODEL_DIR, "lda.model")

lda_ready = os.path.exists(DICT_PATH) and os.path.exists(LDA_PATH)
if lda_ready:
    dictionary = corpora.Dictionary.load(DICT_PATH)
    lda = LdaModel.load(LDA_PATH)

# ----------------------------
# Domain Detection Functions
# ----------------------------
def detect_text_domain(text, noun_tokens):
    """
    Analyze text to determine its domain/category based on vocabulary and patterns.
    """
    text_lower = text.lower()
    words_set = set(nltk.word_tokenize(text_lower))
    
    # Domain indicator keywords with weights
    domain_patterns = {
        'News & Current Affairs': {
            'keywords': ['breaking', 'reported', 'according', 'sources', 'reuters', 'ap', 'cnn', 'bbc', 
                        'journalist', 'correspondent', 'statement', 'official', 'government', 'president',
                        'minister', 'election', 'political', 'parliament', 'congress', 'senate'],
            'weight': 1.0
        },
        'Product Reviews': {
            'keywords': ['review', 'rating', 'stars', 'recommend', 'purchase', 'bought', 'quality',
                        'price', 'value', 'money', 'worth', 'excellent', 'terrible', 'amazing',
                        'disappointed', 'satisfied', 'customer', 'product', 'item', 'delivery'],
            'weight': 1.2
        },
        'Entertainment & Media': {
            'keywords': ['movie', 'film', 'actor', 'actress', 'director', 'cinema', 'theater',
                        'music', 'song', 'album', 'artist', 'band', 'concert', 'performance',
                        'show', 'episode', 'season', 'character', 'plot', 'story', 'entertainment'],
            'weight': 1.1
        },
        'Technology & Science': {
            'keywords': ['technology', 'tech', 'software', 'hardware', 'algorithm', 'data',
                        'research', 'study', 'experiment', 'science', 'scientific', 'innovation',
                        'development', 'programming', 'computer', 'digital', 'artificial', 'intelligence'],
            'weight': 1.0
        },
        'Business & Finance': {
            'keywords': ['business', 'company', 'corporation', 'market', 'stock', 'investment',
                        'profit', 'revenue', 'sales', 'financial', 'economy', 'economic',
                        'bank', 'banking', 'finance', 'money', 'dollar', 'quarter', 'earnings'],
            'weight': 1.0
        },
        'Health & Medical': {
            'keywords': ['health', 'medical', 'doctor', 'patient', 'treatment', 'disease',
                        'medicine', 'therapy', 'hospital', 'clinic', 'symptoms', 'diagnosis',
                        'pharmaceutical', 'drug', 'vaccine', 'healthcare', 'wellness'],
            'weight': 1.1
        },
        'Sports': {
            'keywords': ['game', 'team', 'player', 'match', 'season', 'championship', 'tournament',
                        'score', 'goal', 'win', 'lost', 'victory', 'defeat', 'coach', 'league',
                        'football', 'basketball', 'baseball', 'soccer', 'sports'],
            'weight': 1.2
        },
        'Education & Academic': {
            'keywords': ['university', 'college', 'school', 'student', 'teacher', 'professor',
                        'education', 'academic', 'learning', 'course', 'curriculum', 'degree',
                        'research', 'study', 'knowledge', 'scholarly', 'publication'],
            'weight': 1.0
        },
        'Travel & Tourism': {
            'keywords': ['travel', 'trip', 'vacation', 'hotel', 'restaurant', 'tourist', 'tourism',
                        'destination', 'flight', 'airport', 'booking', 'reservation', 'sightseeing',
                        'culture', 'local', 'experience', 'adventure'],
            'weight': 1.1
        },
        'Lifestyle & Personal': {
            'keywords': ['life', 'personal', 'experience', 'feeling', 'think', 'believe',
                        'opinion', 'lifestyle', 'daily', 'routine', 'family', 'friends',
                        'relationship', 'love', 'happiness', 'advice', 'tips'],
            'weight': 0.9
        }
    }
    
    domain_scores = {}
    
    for domain, config in domain_patterns.items():
        score = 0
        keywords = config['keywords']
        weight = config['weight']
        
        # Count keyword matches
        for keyword in keywords:
            if keyword in text_lower:
                score += text_lower.count(keyword) * weight
        
        # Bonus for exact phrase matches
        phrases = {
            'News & Current Affairs': ['breaking news', 'press release', 'according to sources'],
            'Product Reviews': ['five stars', 'would recommend', 'great value'],
            'Entertainment & Media': ['box office', 'red carpet', 'behind the scenes'],
            'Technology & Science': ['machine learning', 'artificial intelligence', 'peer review'],
            'Business & Finance': ['quarterly earnings', 'market cap', 'stock price'],
            'Sports': ['home team', 'final score', 'world championship']
        }
        
        if domain in phrases:
            for phrase in phrases[domain]:
                if phrase in text_lower:
                    score += 3 * weight
        
        domain_scores[domain] = score
    
    # Normalize scores by text length
    text_length = len(nltk.word_tokenize(text))
    if text_length > 0:
        for domain in domain_scores:
            domain_scores[domain] = domain_scores[domain] / text_length * 100
    
    # Sort domains by score
    sorted_domains = sorted(domain_scores.items(), key=lambda x: x[1], reverse=True)
    
    # Calculate confidence based on score difference
    if len(sorted_domains) > 1:
        top_score = sorted_domains[0][1]
        second_score = sorted_domains[1][1]
        confidence = min(95, max(60, (top_score - second_score) * 10 + 60))
    else:
        confidence = 60
    
    # Return top 5 domains with scores
    return sorted_domains[:5], confidence

def calculate_vocabulary_overlap(noun_tokens, dictionary, min_overlap_ratio=0.3, min_words_in_vocab=10):
    """
    Calculate vocabulary overlap between current text and LDA model's dictionary.
    Returns overlap ratio and recommendation for method selection.
    """
    if not dictionary or not noun_tokens:
        return 0.0, "top_nouns", "No vocabulary or dictionary available"
    
    # Convert tokens to bag of words
    bow = dictionary.doc2bow(noun_tokens)
    words_in_vocab = len(bow)  # Number of words that exist in dictionary
    total_unique_words = len(set(noun_tokens))  # Total unique words in text
    
    if total_unique_words == 0:
        return 0.0, "top_nouns", "No words found"
    
    overlap_ratio = words_in_vocab / total_unique_words
    
    # Decision logic
    if words_in_vocab < min_words_in_vocab:
        return overlap_ratio, "top_nouns", f"Too few vocabulary matches ({words_in_vocab} < {min_words_in_vocab})"
    elif overlap_ratio < min_overlap_ratio:
        return overlap_ratio, "top_nouns", f"Low vocabulary overlap ({overlap_ratio:.1%} < {min_overlap_ratio:.1%})"
    else:
        return overlap_ratio, "lda", f"Good vocabulary overlap ({overlap_ratio:.1%})"

def get_lda_domain_indicators(dictionary):
    """
    Analyze the LDA model's vocabulary to determine what domain it was trained on.
    """
    if not dictionary:
        return "Unknown domain"
    
    # Common word patterns for different domains
    domain_indicators = {
        'news_business': ['reuters', 'company', 'inc', 'corp', 'said', 'wednesday', 'thursday', 'tuesday', 'monday', 'friday'],
        'academic': ['research', 'study', 'analysis', 'theory', 'method', 'paper', 'journal', 'university'],
        'social_media': ['tweet', 'like', 'share', 'follow', 'hashtag', 'post', 'comment'],
        'literature': ['character', 'story', 'novel', 'author', 'book', 'chapter', 'plot'],
        'technical': ['system', 'data', 'algorithm', 'process', 'method', 'implementation', 'framework']
    }
    
    vocab_words = set(dictionary.token2id.keys())
    domain_scores = {}
    
    for domain, indicators in domain_indicators.items():
        score = sum(1 for word in indicators if word in vocab_words)
        domain_scores[domain] = score
    
    if max(domain_scores.values()) > 0:
        likely_domain = max(domain_scores, key=domain_scores.get)
        return likely_domain.replace('_', '/').title()
    else:
        return "General/Mixed domain"

# ----------------------------
# Preprocessing functions
# ----------------------------
def clean_text(text):
    text = text.lower()
    text = re.sub(r"http\S+|www\S+|https\S+", "", text)
    text = re.sub(r"<.*?>", "", text)
    text = re.sub(r"[^a-z\s]", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text

def remove_stopwords(text):
    toks = nltk.word_tokenize(text)
    toks = [w for w in toks if w not in stop_words]
    return " ".join(toks)

def lemmatize_text(text):
    toks = nltk.word_tokenize(text)
    toks = [lemmatizer.lemmatize(w) for w in toks]
    return " ".join(toks)

def nouns_only_tokens(text):
    toks = nltk.word_tokenize(text)
    return [w for w, p in pos_tag(toks) if p.startswith("NN") and len(w) > 1]

def top_nouns(tokens, k=10):
    c = Counter(tokens)
    return c.most_common(k), sum(c.values())

# ----------------------------
# Enhanced Plotting Functions
# ----------------------------
def create_enhanced_bar_chart(x, y, title, x_label, y_label, color_scheme='viridis'):
    fig = go.Figure(data=[
        go.Bar(
            x=x, y=y,
            marker=dict(
                color=y,
                colorscale=color_scheme,
                showscale=True,
                colorbar=dict(thickness=15, len=0.7)
            ),
            text=[f'{val:.2f}' if isinstance(val, float) else str(val) for val in y],
            textposition='outside',
            hovertemplate='<b>%{x}</b><br>Value: %{y}<extra></extra>'
        )
    ])
    
    fig.update_layout(
        title=dict(
            text=title,
            font=dict(size=20, color='white', family='Orbitron'),
            x=0.5
        ),
        xaxis_title=x_label,
        yaxis_title=y_label,
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        font=dict(color='white'),
        xaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
        yaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
        margin=dict(t=60, b=60, l=60, r=60)
    )
    
    return fig

def create_sentiment_gauge(compound_score):
    fig = go.Figure(go.Indicator(
        mode = "gauge+number+delta",
        value = compound_score,
        domain = {'x': [0, 1], 'y': [0, 1]},
        title = {'text': "Sentiment Score", 'font': {'color': 'white', 'size': 20}},
        delta = {'reference': 0, 'increasing': {'color': "RebeccaPurple"}},
        gauge = {
            'axis': {'range': [-1, 1], 'tickcolor': "white"},
            'bar': {'color': "darkblue"},
            'steps': [
                {'range': [-1, -0.05], 'color': "#ff6b6b"},
                {'range': [-0.05, 0.05], 'color': "#ffd93d"},
                {'range': [0.05, 1], 'color': "#4ecdc4"}
            ],
            'threshold': {
                'line': {'color': "red", 'width': 4},
                'thickness': 0.75,
                'value': 0.9
            }
        }
    ))
    
    fig.update_layout(
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        font={'color': "white", 'size': 12},
        height=400
    )
    
    return fig

def create_domain_radar_chart(domain_scores, confidence):
    """Create a radar chart showing domain probabilities"""
    domains = [domain for domain, _ in domain_scores]
    scores = [score for _, score in domain_scores]
    
    fig = go.Figure()
    
    fig.add_trace(go.Scatterpolar(
        r=scores,
        theta=domains,
        fill='toself',
        name='Domain Probability',
        line=dict(color='#00d4ff', width=2),
        fillcolor='rgba(0, 212, 255, 0.3)'
    ))
    
    fig.update_layout(
        polar=dict(
            radialaxis=dict(
                visible=True,
                range=[0, max(scores) * 1.1] if scores else [0, 1],
                gridcolor='rgba(255,255,255,0.2)',
                tickcolor='white'
            ),
            angularaxis=dict(
                gridcolor='rgba(255,255,255,0.2)',
                tickcolor='white'
            )
        ),
        showlegend=True,
        title=dict(
            text=f"🎯 Domain Analysis - {confidence:.0f}% Confidence",
            font=dict(size=20, color='white', family='Orbitron'),
            x=0.5
        ),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        font=dict(color='white'),
        height=500
    )
    
    return fig

# ----------------------------
# Text Analysis Functions
# ----------------------------
def analyze_text_characteristics(text, noun_tokens):
    """Analyze various characteristics of the text"""
    words = nltk.word_tokenize(text)
    sentences = nltk.sent_tokenize(text)
    
    # Basic statistics
    word_count = len(words)
    sentence_count = len(sentences)
    avg_words_per_sentence = word_count / sentence_count if sentence_count > 0 else 0
    
    # Vocabulary richness (unique words / total words)
    unique_words = len(set(word.lower() for word in words if word.isalpha()))
    vocab_richness = unique_words / word_count if word_count > 0 else 0
    
    # Reading complexity (approximate)
    if sentence_count > 0 and word_count > 0:
        # Simple readability estimate
        avg_sentence_length = word_count / sentence_count
        complexity_score = min(100, max(0, (avg_sentence_length - 10) * 5 + 50))
    else:
        complexity_score = 50
    
    # Content density (nouns per total words)
    content_density = len(noun_tokens) / word_count if word_count > 0 else 0
    
    return {
        'word_count': word_count,
        'sentence_count': sentence_count,
        'avg_words_per_sentence': avg_words_per_sentence,
        'vocab_richness': vocab_richness,
        'complexity_score': complexity_score,
        'content_density': content_density,
        'unique_words': unique_words
    }

# ----------------------------
# Main Application Logic
# ----------------------------

# Load custom CSS and create floating particles
load_custom_css()
create_floating_particles()

# Title section
create_title_section()

# Main interface
col1, col2 = st.columns([2, 1])

with col1:
    create_section_header("📝 INPUT CONFIGURATION", "⚙️")
    
    input_type = st.radio(
        'Select Input Method:', 
        ('✍️ Enter Text Manually', '📂 Upload Document'), 
        horizontal=True
    )
    
    # Analysis method selection
    create_section_header("🧠 ANALYSIS METHOD", "🔍")
    topic_method = st.radio(
        'Choose Analysis Method:',
        ('🤖 Smart Auto-Select (Recommended)', '🏷️ Top Keywords Only', '🧠 LDA Only'),
        horizontal=False,
        help="Smart Auto-Select analyzes vocabulary overlap to choose the best method automatically"
    )

with col2:
    create_section_header("📊 ANALYSIS INFO", "📈")
    
    # Show LDA model status and domain
    if lda_ready:
        domain = get_lda_domain_indicators(dictionary)
        vocab_size = len(dictionary.token2id)
        st.markdown(f"""
        <div class="info-card">
            <h4 style="color: #4ecdc4; margin: 0;">✅ LDA Model Ready</h4>
            <p style="color: #a0a0a0; margin: 5px 0;"><strong>Domain:</strong> {domain}</p>
            <p style="color: #a0a0a0; margin: 5px 0;"><strong>Vocabulary:</strong> {vocab_size:,} words</p>
            <p style="color: #a0a0a0; margin: 5px 0;"><strong>Topics:</strong> {lda.num_topics}</p>
        </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown("""
        <div class="info-card">
            <h4 style="color: #ff6b6b; margin: 0;">❌ LDA Model Not Found</h4>
            <p style="color: #a0a0a0; margin: 5px 0;">Will use Domain Detection analysis</p>
        </div>
        """, unsafe_allow_html=True)

# Text input section
text = ''
create_section_header("📄 TEXT INPUT", "📝")

if input_type == '✍️ Enter Text Manually':
    text = st.text_area(
        '✨ Enter your text here for AI analysis:', 
        height=200,
        placeholder="Paste your text here... The AI will automatically detect the domain and select the best analysis method!"
    )
else:
    uploaded_file = st.file_uploader(
        '🚀 Upload your document:', 
        type=['pdf', 'txt', 'docx', 'pptx'],
        help="Supported formats: PDF, TXT, DOCX, PPTX"
    )
    
    if uploaded_file is not None:
        with st.spinner('🔄 Processing your document...'):
            file_type = uploaded_file.type
            try:
                if file_type == 'text/plain':
                    text = uploaded_file.read().decode('utf-8', errors='ignore')
                elif file_type == 'application/pdf':
                    pdf_reader = PyPDF2.PdfReader(uploaded_file)
                    text = "".join([page.extract_text() or "" for page in pdf_reader.pages])
                elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    doc = Document(uploaded_file)
                    text = "\n".join([para.text for para in doc.paragraphs])
                elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                    ppt = Presentation(uploaded_file)
                    text = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
                else:
                    st.error('❌ Unsupported file type.')
            except Exception as e:
                st.error(f'❌ Error reading file: {e}')

# Analysis section
if text:
    # Text preview
    create_section_header("👀 TEXT PREVIEW", "📖")
    preview_text = text[:500] + ('...' if len(text) > 500 else '')
    create_info_card(f"<pre style='white-space: pre-wrap; font-size: 0.9rem;'>{preview_text}</pre>")
    
    # Processing with spinner
    with st.spinner('🧠 AI is analyzing your text... Please wait'):
        time.sleep(1)  # Dramatic pause for effect
        cleaned = clean_text(text)
        no_stop = remove_stopwords(cleaned)
        lemmad = lemmatize_text(no_stop)
        noun_tokens = nouns_only_tokens(lemmad)
    
    # Smart method selection logic
    selected_method = "top_nouns"  # Default fallback
    overlap_ratio = 0.0
    selection_reason = ""
    
    if topic_method == '🤖 Smart Auto-Select (Recommended)' and lda_ready:
        overlap_ratio, selected_method, selection_reason = calculate_vocabulary_overlap(noun_tokens, dictionary)
        
        # Display method selection info
        create_section_header("🎯 SMART METHOD SELECTION", "🤖")
        if selected_method == "lda":
            create_method_alert(f"""
                <strong>🧠 Using LDA Analysis</strong><br>
                <strong>Vocabulary Overlap:</strong> {overlap_ratio:.1%}<br>
                <strong>Reason:</strong> {selection_reason}<br>
                <strong>✅ Good match with training domain</strong>
            """)
        else:
            create_method_alert(f"""
                <strong>🏷️ Using Domain Detection</strong><br>
                <strong>Vocabulary Overlap:</strong> {overlap_ratio:.1%}<br>
                <strong>Reason:</strong> {selection_reason}<br>
                <strong>⚠️ Text domain differs from LDA training data</strong>
            """)
    elif topic_method == '🏷️ Top Keywords Only':
        selected_method = "top_nouns"
        create_method_alert("🏷️ <strong>Using Domain Detection</strong> (User Selected)")
    elif topic_method == '🧠 LDA Only':
        if lda_ready:
            selected_method = "lda"
            create_method_alert("🧠 <strong>Using LDA Analysis</strong> (User Selected)")
        else:
            selected_method = "top_nouns"
            create_method_alert("🏷️ <strong>Fallback to Domain Detection</strong> (LDA model not available)")
    
    # Domain Analysis
    domain_scores, confidence = detect_text_domain(text, noun_tokens)
    text_characteristics = analyze_text_characteristics(text, noun_tokens)
    
    # Results in columns
    col1, col2 = st.columns(2)
    
    with col1:
        # Title Generation
        create_section_header("🏆 GENERATED TITLE", "📰")
        if noun_tokens:
            top1, _ = top_nouns(noun_tokens, k=1)
            title_word = top1[0][0].title()
            create_metric_card("Generated Title", f"📚 {title_word}", "Based on most frequent key noun")
        else:
            fallback = " ".join(nltk.word_tokenize(lemmad)[:10]).capitalize() + "..."
            create_metric_card("Generated Title", f"📚 {fallback}", "Fallback title generated")
        
        # Summary - REMOVED LENGTH RESTRICTION
        create_section_header("📝 SMART SUMMARY", "🔍")
        try:
            with st.spinner('✨ Generating intelligent summary...'):
                # Try different ratios for better summary
                summary = None
                for ratio in [0.3, 0.2, 0.1]:
                    try:
                        summary = summarize(text, ratio=ratio)
                        if summary and len(summary.strip()) > 50:
                            break
                    except:
                        continue
                
                if not summary:
                    # Fallback: create summary from first few sentences
                    sentences = nltk.sent_tokenize(text)
                    if len(sentences) > 2:
                        summary = '. '.join(sentences[:2]) + '.'
                    else:
                        summary = text[:300] + ('...' if len(text) > 300 else '')
                
                create_info_card(f"<div style='font-size: 1.1rem; line-height: 1.6;'>{summary}</div>")
        except Exception as e:
            # Always provide a summary
            sentences = nltk.sent_tokenize(text)
            if len(sentences) > 1:
                summary = '. '.join(sentences[:2]) + '.'
            else:
                summary = text[:300] + ('...' if len(text) > 300 else '')
            create_info_card(f"<div style='font-size: 1.1rem; line-height: 1.6;'>{summary}</div>")
    
    with col2:
        # Sentiment Analysis
        create_section_header("💭 SENTIMENT ANALYSIS", "🎭")
        sentiment = analyzer.polarity_scores(text)
        compound = sentiment['compound']
        
        if compound >= 0.05:
            mood = 'Positive 😊'
            mood_color = 'status-positive'
        elif compound <= -0.05:
            mood = 'Negative 😞'
            mood_color = 'status-negative'
        else:
            mood = 'Neutral 😐'
            mood_color = 'status-neutral'
        
        create_metric_card(
            "Overall Sentiment", 
            f"<span class='{mood_color}'>{mood}</span>",
            f"Confidence Score: {abs(compound):.2f}"
        )
        
        # Sentiment Gauge
        fig_gauge = create_sentiment_gauge(compound)
        st.plotly_chart(fig_gauge, use_container_width=True)
    
    # Domain Analysis (Full Width) - REPLACED TOPIC ANALYSIS
    create_section_header("🌍 DOMAIN DETECTION", "🔍")
    
    if domain_scores:
        primary_domain = domain_scores[0][0]
        primary_score = domain_scores[0][1]
        
        create_metric_card(
            "Detected Domain", 
            f"🎯 {primary_domain}",
            f"Confidence: {confidence:.0f}% | Score: {primary_score:.2f}"
        )
        
        # Domain radar chart
        if len(domain_scores) > 1:
            fig_radar = create_domain_radar_chart(domain_scores, confidence)
            st.plotly_chart(fig_radar, use_container_width=True)
        
        # Top domains bar chart
        domains = [domain for domain, _ in domain_scores]
        scores = [score for _, score in domain_scores]
        
        fig_domains = create_enhanced_bar_chart(
            domains, scores,
            '🎯 Domain Probability Scores',
            'Domains', 'Probability Score',
            'plasma'
        )
        st.plotly_chart(fig_domains, use_container_width=True)
    else:
        create_info_card("⚠️ Unable to detect domain from the provided text.")
    
    # Keyword Analysis - Based on selected method
    if selected_method == "top_nouns":
        create_section_header("🏷️ KEY TERMS ANALYSIS", "🔍")
        if noun_tokens:
            top_items, total = top_nouns(noun_tokens, k=10)
            words = [w for w, c in top_items]
            counts = [c for w, c in top_items]
            
            create_metric_card(
                "Most Frequent Term", 
                f"🏅 {words[0]}", 
                f"Appears {counts[0]} times ({counts[0]/total:.1%} of all key terms)"
            )
            
            fig_noun = create_enhanced_bar_chart(
                words, counts,
                '🏷️ Top Keywords Analysis',
                'Keywords', 'Frequency',
                'viridis'
            )
            st.plotly_chart(fig_noun, use_container_width=True)
        else:
            create_info_card("⚠️ No significant terms found after preprocessing.")
    
    else:
         create_section_header("🧠 LDA TOPIC ANALYSIS", "🔍")
         if not lda_ready:
             st.warning('⚠️ LDA model/dictionary not found. Showing keyword analysis instead.')
             if noun_tokens:
                 top_items, total = top_nouns(noun_tokens, k=10)
                 words = [w for w, c in top_items]
                 counts = [c for w, c in top_items]
                
                 create_metric_card(
                     "Most Frequent Term", 
                     f"🏅 {words[0]}", 
                     f"Using Keywords - {counts[0]} occurrences"
                 )
                
                 fig_noun = create_enhanced_bar_chart(
                     words, counts,
                     '🏷️ Fallback: Keywords Analysis',
                     'Keywords', 'Frequency',
                     'plasma'
                 )
                 st.plotly_chart(fig_noun, use_container_width=True)
         else:
            bow = dictionary.doc2bow(noun_tokens)
            if not bow:
                create_info_card("⚠️ No words from your text match the LDA vocabulary. The model was likely trained on a different domain.")
                domain = get_lda_domain_indicators(dictionary)
                st.info(f"💡 **Tip**: Your LDA model appears to be trained on **{domain}** data. Try text from a similar domain for better results.")
            else:
                topics = lda.get_document_topics(bow)
                if topics:
                    topics_sorted = sorted(topics, key=lambda x: x[1], reverse=True)
                    dominant_topic, dominant_prob = topics_sorted[0]
                    
                    create_metric_card(
                        "Dominant LDA Topic", 
                        f"🧠 Topic {dominant_topic}",
                        f"Probability: {dominant_prob:.2%} | Vocab Overlap: {overlap_ratio:.1%}"
                    )
                    
                    topic_words = lda.show_topic(dominant_topic, 10)
                    words = [w for w, p in topic_words]
                    probs = [p for w, p in topic_words]
                    
                    st.markdown(f"**🔍 Key Terms:** {', '.join(words[:5])}")
                    
                    fig_topic = create_enhanced_bar_chart(
                        words, probs,
                        f'🧠 LDA Topic {dominant_topic} - Word Distribution',
                        'Words', 'Probability',
                        'viridis'
                    )
                    st.plotly_chart(fig_topic, use_container_width=True)
                else:
                    create_info_card("⚠️ No topics detected (text may be too short or vocabulary mismatch).")
    
    # Text Characteristics Analysis
    create_section_header("📊 TEXT CHARACTERISTICS", "📈")
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        create_metric_card("Word Count", f"{text_characteristics['word_count']:,}", "Total words in text")
    with col2:
        create_metric_card("Sentences", f"{text_characteristics['sentence_count']}", "Number of sentences")
    with col3:
        create_metric_card("Vocabulary Richness", f"{text_characteristics['vocab_richness']:.2f}", "Unique words / Total words")
    with col4:
        create_metric_card("Content Density", f"{text_characteristics['content_density']:.2f}", "Key terms / Total words")
    
    # Additional Sentiment Details
    create_section_header("📈 DETAILED SENTIMENT BREAKDOWN", "🔍")
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        create_metric_card("Positive", f"{sentiment['pos']:.2f}", "Joy, happiness, optimism")
    with col2:
        create_metric_card("Neutral", f"{sentiment['neu']:.2f}", "Objective, factual tone")
    with col3:
        create_metric_card("Negative", f"{sentiment['neg']:.2f}", "Sadness, anger, criticism")
    with col4:
        create_metric_card("Compound", f"{sentiment['compound']:.2f}", "Overall sentiment score")
    
    # Sentiment breakdown chart
    scores = [sentiment['pos'], sentiment['neu'], sentiment['neg']]
    labels = ['Positive 😊', 'Neutral 😐', 'Negative 😞']
    colors = ['#4ecdc4', '#ffd93d', '#ff6b6b']
    
    fig_sentiment = go.Figure(data=[
        go.Bar(x=labels, y=scores, marker_color=colors,
              text=[f'{score:.2f}' for score in scores],
              textposition='outside')
    ])
    
    fig_sentiment.update_layout(
        title=dict(text='🎭 Sentiment Component Analysis', 
                  font=dict(size=20, color='white', family='Orbitron'), x=0.5),
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        font=dict(color='white'),
        xaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
        yaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
        showlegend=False
    )
    
    st.plotly_chart(fig_sentiment, use_container_width=True)
    
    # Analysis Summary and Recommendations - FIXED SECTION
    create_section_header("💡 ANALYSIS INSIGHTS", "🔍")
    
    col1, col2 = st.columns(2)
    
    with col1:  # Assuming col1 is a Streamlit column defined earlier
        if domain_scores:
            primary_domain = domain_scores[0][0]
            st.markdown(
            """
            <div class="info-card" style="background-color: #1a1a1a; padding: 15px; border-radius: 8px;">
                <h4 style="color: #4ecdc4; margin: 0;">📊 Analysis Summary</h4>
            </div>
            <div>    
                <p style="margin: 5px 0;">🎯 <strong>Primary Domain:</strong> {}</p>
                <p style="margin: 5px 0;">📈 <strong>Confidence Level:</strong> {:.0f}%</p>
                <p style="margin: 5px 0;">💭 <strong>Sentiment:</strong> {}</p>
                <p style="margin: 5px 0;">📝 <strong>Complexity:</strong> {}</p>
            </div>
            """.format(
                primary_domain,
                confidence,
                mood,
                "High" if text_characteristics["complexity_score"] > 70 else "Medium" if text_characteristics["complexity_score"] > 40 else "Low"
            ),
            unsafe_allow_html=True
            )

        else:
            st.markdown(
            """
            <div class="info-card" style="background-color: #1a1a1a; padding: 15px; border-radius: 8px;">
                <h4 style="color: #4ecdc4; margin: 0;">📊 Analysis Summary</h4>
                <p style="color: #a0a0a0; margin: 5px 0;">No analysis results available.</p>
            </div>
            """,
            unsafe_allow_html=True
        )
    
    with col2:
        st.markdown("""
        <div class="info-card">
            <h4 style="color: #ff6b6b; margin: 0;">💡 Key Insights</h4>
        """, unsafe_allow_html=True)
        
        # Provide meaningful insights based on analysis
        insights = []
        
        # Domain-based insights
        if domain_scores:
            primary_domain = domain_scores[0][0]
            if confidence > 80:
                insights.append(f"✅ Strong {primary_domain.lower()} content detected")
            elif confidence > 60:
                insights.append(f"📊 Likely {primary_domain.lower()} content")
            else:
                insights.append("🔄 Mixed domain content detected")
        
        # Sentiment insights
        if abs(compound) > 0.3:
            insights.append(f"{'📈' if compound > 0 else '📉'} Strong emotional tone detected")
        else:
            insights.append("⚖️ Balanced, neutral tone")
        
        # Text quality insights
        if text_characteristics['vocab_richness'] > 0.6:
            insights.append("📚 Rich vocabulary usage")
        elif text_characteristics['vocab_richness'] < 0.4:
            insights.append("🔄 Consider diversifying vocabulary")
        
        # Content density insights
        if text_characteristics['content_density'] > 0.3:
            insights.append("💎 High information density")
        else:
            insights.append("💬 Conversational style content")
        
        for insight in insights:
            st.markdown(f"<p style='color: #a0a0a0; margin: 5px 0;'>{insight}</p>", unsafe_allow_html=True)
        
        st.markdown("</div>", unsafe_allow_html=True)
    
    # Additional Analysis Recommendations
    create_section_header("🚀 OPTIMIZATION SUGGESTIONS", "💡")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown("""
        <div class="info-card">
            <h4 style="color: #4ecdc4; margin: 0;">🎯 Content Optimization</h4>
        """, unsafe_allow_html=True)
        
        if text_characteristics['avg_words_per_sentence'] > 25:
            st.markdown("<p style='color: #a0a0a0; margin: 5px 0;'>✂️ Consider shorter sentences for better readability</p>", unsafe_allow_html=True)
        elif text_characteristics['avg_words_per_sentence'] < 10:
            st.markdown("<p style='color: #a0a0a0; margin: 5px 0;'>🔗 Consider combining short sentences for flow</p>", unsafe_allow_html=True)
        else:
            st.markdown("<p style='color: #a0a0a0; margin: 5px 0;'>✅ Good sentence length balance</p>", unsafe_allow_html=True)
        
        st.markdown("</div>", unsafe_allow_html=True)

        create_section_header("📥 DOWNLOAD YOUR REPORT", "📄")
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        story.append(Paragraph("Narrative Nexus - AI Text Analysis Report", styles['Title']))
        story.append(Spacer(1, 12))
        story.append(Paragraph(f"<b>Generated Title:</b> {title_word}", styles['Normal']))
        story.append(Paragraph(f"<b>Detected Domain:</b> {primary_domain}", styles['Normal']))
        story.append(Paragraph(f"<b>Sentiment:</b> {mood}", styles['Normal']))
        story.append(Paragraph(f"<b>Confidence:</b> {confidence:.0f}%", styles['Normal']))
        story.append(Spacer(1, 12))
        story.append(Paragraph("<b>Summary:</b>", styles['Heading2']))
        story.append(Paragraph(summary, styles['Normal']))
        story.append(Spacer(1, 12))
        story.append(Paragraph("<b>Insights:</b>", styles['Heading2']))
        for insight in insights:
            story.append(Paragraph(insight, styles['Normal']))

        doc.build(story)
        pdf_data = buffer.getvalue()

        st.download_button(
            label="📥 Download Full Report (PDF)",
            data=pdf_data,
            file_name="narrative_nexus_report.pdf",
            mime="application/pdf",
        )
    
    with col2:
        st.markdown("""
        <div class="info-card">
            <h4 style="color: #ffd93d; margin: 0;">📊 Audience Targeting</h4>
        """, unsafe_allow_html=True)
        
        if domain_scores:
            primary_domain = domain_scores[0][0]
            domain_suggestions = {
                'News & Current Affairs': 'Focus on timeliness and credibility',
                'Product Reviews': 'Emphasize pros/cons and value proposition',
                'Entertainment & Media': 'Highlight emotional engagement',
                'Technology & Science': 'Ensure technical accuracy',
                'Business & Finance': 'Include relevant metrics and data',
                'Health & Medical': 'Prioritize accuracy and safety',
                'Sports': 'Add performance statistics',
                'Education & Academic': 'Structure for learning outcomes',
                'Travel & Tourism': 'Include practical information',
                'Lifestyle & Personal': 'Focus on relatability'
            }
            
            suggestion = domain_suggestions.get(primary_domain, 'Tailor content to target audience')
            st.markdown(f"<p style='color: #a0a0a0; margin: 5px 0;'>💡 {suggestion}</p>", unsafe_allow_html=True)
        
        st.markdown("</div>", unsafe_allow_html=True)
    
    with col3:
        st.markdown("""
        <div class="info-card">
            <h4 style="color: #ff6b6b; margin: 0;">🎭 Tone Adjustment</h4>
        """, unsafe_allow_html=True)
        
        if compound > 0.5:
            st.markdown("<p style='color: #a0a0a0; margin: 5px 0;'>😊 Very positive tone - great for engagement</p>", unsafe_allow_html=True)
        elif compound < -0.5:
            st.markdown("<p style='color: #a0a0a0; margin: 5px 0;'>⚠️ Very negative tone - consider balance</p>", unsafe_allow_html=True)
        elif abs(compound) < 0.1:
            st.markdown("<p style='color: #a0a0a0; margin: 5px 0;'>📊 Neutral tone - good for informational content</p>", unsafe_allow_html=True)
        else:
            st.markdown("<p style='color: #a0a0a0; margin: 5px 0;'>⚖️ Balanced emotional tone</p>", unsafe_allow_html=True)
        
        st.markdown("</div>", unsafe_allow_html=True)

    
else:
    # Welcome message when no text
    st.markdown("""
    <div style="text-align: center; padding: 60px 20px;">
        <h2 style="color: #00d4ff; font-family: 'Orbitron', monospace;">
            🚀 Ready for Smart Text Analysis!
        </h2>
        <p style="font-size: 1.2rem; color: #a0a0a0; margin: 20px 0;">
            Upload a document or enter text above. The AI will automatically detect the domain and choose the best analysis method.
        </p>
        <div style="font-size: 3rem; margin: 20px 0;">🤖✨📊</div>
        <div style="background: rgba(255, 255, 255, 0.05); border-radius: 15px; padding: 20px; margin: 20px 0;">
            <h3 style="color: #4ecdc4; margin: 0 0 15px 0;">🧠 Smart Features</h3>
            <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 15px; text-align: left;">
                <div>
                    <strong style="color: #00d4ff;">🎯 Domain Detection</strong><br>
                    <span style="color: #a0a0a0; font-size: 0.9rem;">Identifies content type: News, Reviews, Entertainment, etc.</span>
                </div>
                <div>
                    <strong style="color: #00d4ff;">📊 Smart Analysis</strong><br>
                    <span style="color: #a0a0a0; font-size: 0.9rem;">Chooses optimal method based on content characteristics</span>
                </div>
                <div>
                    <strong style="color: #00d4ff;">🔍 Content Insights</strong><br>
                    <span style="color: #a0a0a0; font-size: 0.9rem;">Analyzes readability, complexity, and engagement factors</span>
                </div>
                <div>
                    <strong style="color: #00d4ff;">💡 Actionable Recommendations</strong><br>
                    <span style="color: #a0a0a0; font-size: 0.9rem;">Provides specific suggestions for content optimization</span>
                </div>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)